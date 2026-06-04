"""
Extractors for Microsoft Office documents (DOCX, PPTX, XLSX).
"""

from pathlib import Path
from docx import Document
from pptx import Presentation
import openpyxl
from types import ExtractedContent, DocumentMetadata, FileType
from extractors.base_extractor import BaseExtractor, ExtractionError
from utils.logger import get_logger

logger = get_logger(__name__)


class DOCXExtractor(BaseExtractor):
    """Extractor for .docx files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is DOCX."""
        return file_path.suffix.lower() in [".doc", ".docx"]
    
    async def extract(self, file_path: Path, metadata: DocumentMetadata) -> ExtractedContent:
        """
        Extract content from DOCX file.
        
        Args:
            file_path: Path to DOCX file
            metadata: Document metadata
            
        Returns:
            Extracted content
        """
        try:
            doc = Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_content = [cell.text for cell in row.cells]
                    text_content.append(" | ".join(row_content))
            
            text = "\n".join(text_content)
            
            logger.info(f"Extracted DOCX from {file_path.name}", extra={
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables)
            })
            
            return ExtractedContent(
                metadata=metadata,
                raw_text=text,
                has_tables=len(doc.tables) > 0
            )
            
        except Exception as e:
            logger.error(f"Failed to extract DOCX {file_path}: {str(e)}")
            raise ExtractionError(f"DOCX extraction failed: {str(e)}")


class PPTXExtractor(BaseExtractor):
    """Extractor for .pptx files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is PPTX."""
        return file_path.suffix.lower() in [".ppt", ".pptx"]
    
    async def extract(self, file_path: Path, metadata: DocumentMetadata) -> ExtractedContent:
        """
        Extract content from PPTX file.
        
        Args:
            file_path: Path to PPTX file
            metadata: Document metadata
            
        Returns:
            Extracted content
        """
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"\n--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_content.append(shape.text)
            
            text = "\n".join(text_content)
            
            logger.info(f"Extracted PPTX from {file_path.name}", extra={
                "slides": len(prs.slides)
            })
            
            return ExtractedContent(
                metadata=metadata,
                raw_text=text,
                pages=len(prs.slides)
            )
            
        except Exception as e:
            logger.error(f"Failed to extract PPTX {file_path}: {str(e)}")
            raise ExtractionError(f"PPTX extraction failed: {str(e)}")


class XLSXExtractor(BaseExtractor):
    """Extractor for .xlsx files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is XLSX."""
        return file_path.suffix.lower() in [".xls", ".xlsx"]
    
    async def extract(self, file_path: Path, metadata: DocumentMetadata) -> ExtractedContent:
        """
        Extract content from XLSX file.
        
        Args:
            file_path: Path to XLSX file
            metadata: Document metadata
            
        Returns:
            Extracted content
        """
        try:
            workbook = openpyxl.load_workbook(file_path)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"\n--- Sheet: {sheet_name} ---")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    text_content.append(" | ".join(row_text))
            
            text = "\n".join(text_content)
            
            logger.info(f"Extracted XLSX from {file_path.name}", extra={
                "sheets": len(workbook.sheetnames)
            })
            
            return ExtractedContent(
                metadata=metadata,
                raw_text=text,
                has_tables=True
            )
            
        except Exception as e:
            logger.error(f"Failed to extract XLSX {file_path}: {str(e)}")
            raise ExtractionError(f"XLSX extraction failed: {str(e)}")
